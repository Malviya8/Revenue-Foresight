"""
Problem + Objective slide — Canva Black / Orange / White pitch theme.
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

BLACK = RGBColor(0x0A, 0x0A, 0x0A)
NEAR_BLACK = RGBColor(0x14, 0x14, 0x14)
CARD = RGBColor(0x1A, 0x1A, 0x1A)
ORANGE = RGBColor(0xFF, 0x5C, 0x2A)
ORANGE_DEEP = RGBColor(0xE0, 0x3E, 0x12)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
MUTED = RGBColor(0xA8, 0xA8, 0xA8)
LINE = RGBColor(0x2E, 0x2E, 0x2E)

W, H = Inches(13.333), Inches(7.5)
OUT = Path(__file__).resolve().parent / "Problem_Objective_Slide.pptx"
FONT = "Arial"


def _run(run, text, size=18, bold=False, color=WHITE, italic=False):
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = FONT


def txt(slide, left, top, width, height, lines, *, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    if isinstance(lines, str):
        lines = [lines]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(3)
        if isinstance(line, tuple):
            text, size, bold, color = line[:4]
            italic = line[4] if len(line) > 4 else False
        else:
            text, size, bold, color, italic = line, 16, False, WHITE, False
        r = p.add_run()
        _run(r, text, size=size, bold=bold, color=color, italic=italic)
    return box


def bg(slide, color=BLACK):
    f = slide.background.fill
    f.solid()
    f.fore_color.rgb = color


def shape(slide, kind, left, top, width, height, fill, line=None):
    sh = slide.shapes.add_shape(kind, left, top, width, height)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(1)
    return sh


def rect(slide, left, top, width, height, fill, line=None):
    return shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, width, height, fill, line)


def roundr(slide, left, top, width, height, fill, line=None):
    return shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height, fill, line)


def oval(slide, left, top, width, height, fill):
    return shape(slide, MSO_AUTO_SHAPE_TYPE.OVAL, left, top, width, height, fill)


def main():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)

    # Ambient Canva orbs
    oval(s, Inches(-1.1), Inches(-1.2), Inches(3.2), Inches(3.2), ORANGE)
    oval(s, Inches(11.4), Inches(5.5), Inches(2.8), Inches(2.8), ORANGE_DEEP)

    # Meta
    txt(s, Inches(0.55), Inches(0.28), Inches(5), Inches(0.3), [("REVENUE FORESIGHT", 11, True, MUTED)])
    txt(s, Inches(9.5), Inches(0.28), Inches(3.2), Inches(0.3), [("PROBLEM  ·  OBJECTIVE", 11, True, ORANGE)], align=PP_ALIGN.RIGHT)

    # Title
    txt(s, Inches(0.55), Inches(0.7), Inches(12), Inches(0.85), [("The gap we close.", 40, True, WHITE)])
    txt(
        s,
        Inches(0.55),
        Inches(1.45),
        Inches(12),
        Inches(0.4),
        [("From rear-view dashboards to decision-grade planning.", 15, False, MUTED)],
    )

    # ── LEFT: PROBLEM ──
    roundr(s, Inches(0.55), Inches(2.1), Inches(6.0), Inches(4.7), CARD)
    rect(s, Inches(0.55), Inches(2.1), Inches(0.12), Inches(4.7), ORANGE)

    roundr(s, Inches(0.9), Inches(2.35), Inches(1.7), Inches(0.38), ORANGE)
    txt(s, Inches(0.9), Inches(2.4), Inches(1.7), Inches(0.3), [("PROBLEM", 12, True, WHITE)], align=PP_ALIGN.CENTER)

    txt(
        s,
        Inches(0.9),
        Inches(2.95),
        Inches(5.3),
        Inches(1.1),
        [
            ("Agencies have rich past performance —", 15, False, MUTED),
            ("but weak forward answers under real budget choices.", 15, True, WHITE),
        ],
    )

    pains = [
        ("Fragmented truth", "Google, Meta, Bing speak different data languages."),
        ("Rear-view tools", "Dashboards explain yesterday; planners need 30 / 60 / 90 days."),
        ("False certainty", "A single point forecast hides downside and upside risk."),
        ("Linear spend myth", "“+20% budget ≠ +20% revenue” — yet decisions often assume it."),
    ]
    for i, (t, b) in enumerate(pains):
        top = Inches(4.15) + i * Inches(0.6)
        oval(s, Inches(0.95), top + Inches(0.08), Inches(0.18), Inches(0.18), ORANGE)
        txt(s, Inches(1.3), top, Inches(4.9), Inches(0.25), [(t, 13, True, WHITE)])
        txt(s, Inches(1.3), top + Inches(0.25), Inches(4.9), Inches(0.28), [(b, 11, False, MUTED)])

    # ── RIGHT: OBJECTIVE ──
    roundr(s, Inches(6.8), Inches(2.1), Inches(6.0), Inches(4.7), NEAR_BLACK)
    rect(s, Inches(6.8), Inches(2.1), Inches(6.0), Inches(0.1), ORANGE)

    roundr(s, Inches(7.15), Inches(2.4), Inches(2.0), Inches(0.38), ORANGE)
    txt(s, Inches(7.15), Inches(2.45), Inches(2.0), Inches(0.3), [("OBJECTIVE", 12, True, WHITE)], align=PP_ALIGN.CENTER)

    txt(
        s,
        Inches(7.15),
        Inches(3.0),
        Inches(5.3),
        Inches(1.35),
        [
            ("Deliver probabilistic revenue & ROAS forecasts", 16, True, WHITE),
            ("for the next 30 / 60 / 90 days —", 15, False, MUTED),
            ("conditioned on planned spend, across hierarchy.", 15, False, MUTED),
        ],
    )

    # Outcome chips
    outcomes = [
        ("Ranges", "P10 · P50 · P90"),
        ("Levels", "Store → Channel → Type → Campaign"),
        ("Scenarios", "Budget what-ifs before spend"),
        ("Offline-ready", "One command. No retrain at score time."),
    ]
    for i, (t, b) in enumerate(outcomes):
        col, row = i % 2, i // 2
        left = Inches(7.15) + col * Inches(2.65)
        top = Inches(4.55) + row * Inches(0.95)
        roundr(s, left, top, Inches(2.5), Inches(0.85), CARD)
        txt(s, left + Inches(0.15), top + Inches(0.12), Inches(2.2), Inches(0.3), [(t, 13, True, ORANGE)])
        txt(s, left + Inches(0.15), top + Inches(0.42), Inches(2.2), Inches(0.35), [(b, 11, False, MUTED)])

    # Speaker notes
    s.notes_slide.notes_text_frame.text = (
        "PROBLEM: Agencies have rich past performance but weak forward answers under real budget "
        "choices. Data is fragmented across Google, Meta, and Bing. Dashboards explain yesterday; "
        "planners need 30/60/90-day outlooks. Single point forecasts hide risk. Linear spend "
        "assumptions burn budget.\n\n"
        "OBJECTIVE: Deliver probabilistic revenue and ROAS forecasts for 30/60/90 days, conditioned "
        "on planned spend, at aggregate / channel / type / campaign levels — with ranges, hierarchy "
        "coherence, and budget what-ifs — offline-ready for scoring."
    )

    prs.save(str(OUT))
    print(f"Wrote {OUT}")


if __name__ == "__main__":
    main()
